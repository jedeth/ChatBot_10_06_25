from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import logging

class DocumentProcessor:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.extractors = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.xlsx': self._extract_excel,
            '.pptx': self._extract_pptx,
            '.txt': self._extract_txt,
        }

    def extract_text(self, file_path: str) -> str:
        file_ext = Path(file_path).suffix.lower()
        extractor = self.extractors.get(file_ext)
        if extractor:
            try:
                self.logger.info(f"Extraction du texte du fichier : {file_path}")
                return extractor(file_path)
            except Exception as e:
                self.logger.error(f"Erreur lors de l'extraction de {file_path}: {e}")
                raise ValueError(f"Impossible de traiter le fichier : {e}")
        self.logger.warning(f"Type de fichier non supporté : {file_ext}")
        raise ValueError(f"Type de fichier non supporté : {file_ext}")

    def _extract_pdf(self, file_path: str) -> str:
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text

    def _extract_docx(self, file_path: str) -> str:
        doc = DocxDocument(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_excel(self, file_path: str) -> str:
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Feuille: {sheet_name}\n"
            for row in sheet.iter_rows():
                text += "\t".join([str(cell.value or "") for cell in row]) + "\n"
        return text

    def _extract_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()